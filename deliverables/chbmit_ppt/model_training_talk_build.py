from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(r"E:\CHB-MIT\deliverables\chbmit_ppt\Model_Training_Talk_EN.pptx")

BG = RGBColor(0xF6, 0xF4, 0xEF)
INK = RGBColor(0x17, 0x32, 0x4D)
BODY = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x64, 0x74, 0x8B)
LINE = RGBColor(0xD6, 0xDE, 0xE8)
TEAL = RGBColor(0x1F, 0x7A, 0x52)
TEAL_FILL = RGBColor(0xE8, 0xF6, 0xEE)
BLUE = RGBColor(0x1F, 0x6A, 0xA5)
BLUE_FILL = RGBColor(0xEA, 0xF3, 0xFB)
AMBER = RGBColor(0xC9, 0x8B, 0x28)
AMBER_FILL = RGBColor(0xFF, 0xF4, 0xDE)
RED = RGBColor(0xC8, 0x4C, 0x3A)
RED_FILL = RGBColor(0xFB, 0xEA, 0xEA)
SLATE_FILL = RGBColor(0xEE, 0xF2, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    top = slide.shapes.add_shape(SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.34))
    top.fill.solid()
    top.fill.fore_color.rgb = INK
    top.line.color.rgb = INK


def add_text(slide, x, y, w, h, text, size=14, color=BODY, bold=False, name="Aptos", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    box.line.color.rgb = BG
    box.fill.background()
    return box


def add_panel(slide, x, y, w, h, fill=WHITE, line=LINE):
    shp = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    return shp


def add_header(slide, title, subtitle):
    set_bg(slide)
    add_text(slide, 0.42, 0.48, 11.2, 0.45, title, size=25, color=INK, bold=True, name="Aptos Display")
    add_text(slide, 0.42, 0.92, 11.8, 0.2, subtitle, size=10, color=MUTED)


def add_footer(slide, page):
    line = slide.shapes.add_shape(SHAPE.RECTANGLE, Inches(0.42), Inches(7.16), Inches(12.49), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.color.rgb = LINE
    add_text(slide, 0.42, 7.19, 2.8, 0.14, "CHB-MIT | Model & Training", size=8, color=MUTED)
    add_text(slide, 12.45, 7.18, 0.3, 0.14, f"{page:02d}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_bullet_block(slide, x, y, lines, width=2.9, gap=0.46, size=13):
    yy = y
    for line in lines:
        dot = slide.shapes.add_shape(SHAPE.OVAL, Inches(x), Inches(yy + 0.07), Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = TEAL
        dot.line.color.rgb = TEAL
        add_text(slide, x + 0.14, yy, width, 0.32, line, size=size, color=BODY)
        yy += gap


def add_stat(slide, x, y, label, value, fill):
    add_panel(slide, x, y, 1.42, 0.72, fill=fill, line=fill)
    add_text(slide, x + 0.12, y + 0.09, 1.05, 0.14, label, size=8.5, color=INK)
    add_text(slide, x + 0.12, y + 0.26, 1.05, 0.26, value, size=18, color=INK, bold=True, name="Aptos Display")


def add_flow_box(slide, x, y, w, title, subtitle, fill):
    add_panel(slide, x, y, w, 1.25, fill=fill, line=fill)
    add_text(slide, x + 0.08, y + 0.12, w - 0.16, 0.22, title, size=15, color=INK, bold=True, name="Aptos Display", align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.1, y + 0.46, w - 0.2, 0.46, subtitle, size=9.5, color=BODY, align=PP_ALIGN.CENTER)


def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Model Design: Why CNN + LSTM Was the Right Choice",
        "My role focused on model architecture and the training pipeline for future seizure-risk prediction.",
    )

    add_stat(slide, 0.62, 1.22, "Input", "18 x 30s", BLUE_FILL)
    add_stat(slide, 2.12, 1.22, "Context", "6 windows", TEAL_FILL)
    add_stat(slide, 3.62, 1.22, "Forecast", "30 min", AMBER_FILL)

    add_panel(slide, 0.58, 2.1, 3.45, 4.55)
    add_text(slide, 0.8, 2.34, 2.2, 0.24, "Design rationale", size=18, color=INK, bold=True, name="Aptos Display")
    add_bullet_block(
        slide,
        0.82,
        2.76,
        [
            "Target: future seizure risk, not seizure detection.",
            "A 1D CNN learns local EEG morphology from each 30 s window.",
            "An LSTM adds 3 minutes of temporal context across windows.",
            "The head outputs a probability, which we later calibrate.",
        ],
        width=2.88,
        gap=0.54,
        size=11.8,
    )

    add_panel(slide, 4.28, 2.1, 8.45, 4.55)
    add_text(slide, 4.5, 2.34, 2.4, 0.24, "Architecture flow", size=18, color=INK, bold=True, name="Aptos Display")

    add_flow_box(slide, 4.56, 3.0, 1.78, "Input sequence", "6 consecutive windows\n18 channels, 256 Hz", SLATE_FILL)
    add_flow_box(slide, 6.98, 3.0, 1.78, "1D CNN encoder", "Learns waveform and spectral\npatterns per window", BLUE_FILL)
    add_flow_box(slide, 9.4, 3.0, 1.65, "LSTM", "Models temporal progression\nacross windows", TEAL_FILL)
    add_flow_box(slide, 11.62, 3.0, 0.82, "p(y=1)", "Calibrated\nfuture-risk\nprobability", AMBER_FILL)

    for x in (6.45, 8.87, 11.12):
        arrow = slide.shapes.add_shape(SHAPE.CHEVRON, Inches(x), Inches(3.37), Inches(0.42), Inches(0.48))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = LINE
        arrow.line.color.rgb = LINE

    add_panel(slide, 4.58, 4.86, 7.65, 1.16, fill=TEAL_FILL, line=TEAL)
    add_text(slide, 4.78, 5.08, 1.95, 0.18, "Final design decision", size=14, color=TEAL, bold=True, name="Aptos Display")
    add_text(
        slide,
        4.78,
        5.35,
        7.0,
        0.38,
        "Use the tuned CNN + LSTM as the main model: it captures temporal evolution better than a pure CNN, while remaining trainable on our GPU and dataset size.",
        size=12,
        color=BODY,
    )
    add_footer(slide, 1)


def build_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Training Strategy: What Worked, What Failed, and What We Selected",
        "I focused on solving split bias, class imbalance, and probability reliability rather than only chasing a bigger model.",
    )

    add_panel(slide, 0.58, 1.22, 3.7, 5.88)
    add_text(slide, 0.82, 1.48, 2.5, 0.24, "Training decisions", size=18, color=INK, bold=True, name="Aptos Display")
    add_bullet_block(
        slide,
        0.84,
        1.93,
        [
            "Rebuilt the patient split after finding a train/val ratio mismatch.",
            "Used real class distribution plus pos_weight, not synthetic balancing.",
            "Used early stopping, LR scheduling, memmap loading, and W&B.",
            "Compared CNN, CNN+LSTM, and Deep Ensemble on held-out patients.",
            "Calibrated the final model: test ECE fell from 0.0497 to 0.0133.",
        ],
        width=3.0,
        gap=0.49,
        size=11.2,
    )
    add_panel(slide, 0.82, 5.68, 3.16, 1.02, fill=AMBER_FILL, line=AMBER)
    add_text(slide, 1.0, 5.9, 1.7, 0.16, "Key technical insight", size=13, color=AMBER, bold=True, name="Aptos Display")
    add_text(
        slide,
        1.0,
        6.16,
        2.72,
        0.32,
        "The highest score did not come from the most complex pipeline. A well-tuned, calibrated CNN + LSTM generalized best.",
        size=10.6,
        color=BODY,
    )

    add_panel(slide, 4.48, 1.22, 8.28, 5.88)
    add_text(slide, 4.72, 1.48, 2.8, 0.24, "Held-out test comparison", size=18, color=INK, bold=True, name="Aptos Display")

    rows = [
        ["Model", "Test AUROC", "AUPRC", "F1", "ECE", "Reading"],
        ["1D CNN baseline", "0.701", "0.118", "0.070", "-", "Stable baseline"],
        ["CNN + LSTM (seq10)", "0.748", "0.123", "0.172", "-", "More temporal context"],
        ["Tuned CNN + LSTM", "0.797", "0.189", "0.189", "0.013", "Selected final model"],
        ["Deep Ensemble", "0.787", "0.167", "0.020", "-", "Explored, not retained"],
    ]
    tbl = slide.shapes.add_table(5, 6, Inches(4.72), Inches(2.0), Inches(7.62), Inches(2.76)).table
    widths = [2.15, 1.15, 0.86, 0.82, 0.72, 1.92]
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)

    for r in range(5):
        for c in range(6):
            cell = tbl.cell(r, c)
            cell.text = rows[r][c]
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            fill = cell.fill
            fill.solid()
            if r == 0:
                fill.fore_color.rgb = SLATE_FILL
            elif r == 3:
                fill.fore_color.rgb = TEAL_FILL
            elif r == 4:
                fill.fore_color.rgb = RED_FILL
            elif r == 1 and c == 3:
                fill.fore_color.rgb = RED_FILL
            else:
                fill.fore_color.rgb = WHITE

            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = "Aptos Display" if r == 0 else "Aptos"
                    run.font.size = Pt(10.5)
                    run.font.bold = r == 0 or r == 3 or (r == 4 and c == 5)
                    if r == 3 and c > 0:
                        run.font.color.rgb = TEAL
                    elif r == 4 and c > 0:
                        run.font.color.rgb = RED
                    elif r == 1 and c == 3:
                        run.font.color.rgb = RED
                    else:
                        run.font.color.rgb = INK if r == 0 else BODY

    add_panel(slide, 4.72, 5.16, 7.6, 1.02, fill=BLUE_FILL, line=BLUE)
    add_text(slide, 4.94, 5.36, 2.2, 0.16, "What I would say out loud", size=13, color=BLUE, bold=True, name="Aptos Display")
    add_text(
        slide,
        4.94,
        5.63,
        7.0,
        0.32,
        "We trained and compared multiple candidates, including a Deep Ensemble. The best final choice was still the tuned, calibrated CNN + LSTM because it achieved the strongest held-out balance of ranking quality, F1, and probability reliability.",
        size=11.1,
        color=BODY,
    )
    add_footer(slide, 2)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build_slide_1(prs)
    build_slide_2(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved {OUT}")


if __name__ == "__main__":
    build()
